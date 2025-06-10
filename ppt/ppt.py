from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat
import requests
from io import BytesIO

def set_slide_background(slide, rgb_color):
    """Sets the slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def set_slide_background_gradient(slide, color1, color2):
    """Sets a gradient background for the slide."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_decorative_shapes(slide):
    """Adds decorative shapes and background elements."""
    # Add subtle geometric shapes in the background
    # Top right corner accent
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8.5), Inches(0), 
        Inches(1.5), Inches(1.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0x6A, 0x0D, 0xAD)  # Purple accent
    shape1.fill.transparency = 0.8  # Make it subtle
    shape1.line.fill.background()
    
    # Bottom left accent
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(-0.5), Inches(6.5), 
        Inches(2), Inches(2)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0x9B, 0x59, 0xB6)  # Lighter purple
    shape2.fill.transparency = 0.9
    shape2.line.fill.background()

def add_icon_shape(slide, icon_type, x, y, size=0.5):
    """Adds simple icon shapes to slides."""
    if icon_type == "check":
        # Create a checkmark using a circle
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0x2E, 0xCC, 0x71)  # Green
        icon.line.color.rgb = RGBColor(0x27, 0xAE, 0x60)
        icon.line.width = Pt(2)
        return icon
    elif icon_type == "star":
        icon = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0xF1, 0xC4, 0x0F)  # Gold
        icon.line.color.rgb = RGBColor(0xD4, 0xAC, 0x0D)
        return icon
    elif icon_type == "gear":
        icon = slide.shapes.add_shape(MSO_SHAPE.GEAR_6, Inches(x), Inches(y), Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0x3D, 0xBD, 0xE0)  # Blue
        icon.line.color.rgb = RGBColor(0x2E, 0x86, 0xAB)
        return icon
    elif icon_type == "lightbulb":
        # Use a diamond shape as lightbulb
        icon = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0xFF, 0xD7, 0x00)  # Bright yellow
        icon.line.color.rgb = RGBColor(0xCC, 0xAC, 0x00)
        return icon
    else:
        # Default to a circle
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(0x3D, 0xBD, 0xE0)
        return icon
    
def add_title_background_banner(slide, title_text):
    """Adds a styled background banner for titles."""
    # Create a background rectangle for the title
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(1.2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x6A, 0x0D, 0xAD)  # Purple
    banner.fill.transparency = 0.1
    banner.line.fill.background()
    
    # Send banner to back so text appears on top
    banner.element.getparent().remove(banner.element)
    slide.shapes._spTree.insert(2, banner.element)

def add_footer(slide, text):
    """Adds a consistent footer to each slide with enhanced styling."""
    # Add footer background
    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8),
        Inches(10), Inches(0.7)
    )
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)  # Dark blue-gray
    footer_bg.line.fill.background()
    
    # Add footer text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
    p.font.size = Pt(10)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Add company logo placeholder (as a shape)
    logo_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(7.0),
        Inches(1), Inches(0.4)
    )
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = RGBColor(0x6A, 0x0D, 0xAD)
    logo_shape.line.fill.background()

def create_presentation(data):
    """
    Generates a PowerPoint presentation based on the provided JSON data,
    mimicking the style of the user-provided template.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Define Color & Font Styles from Template ---
    # Using a professional palette inspired by the template
    BG_COLOR = RGBColor(0xFA, 0xFA, 0xFA) # Off-white background
    TEXT_COLOR = RGBColor(0x33, 0x33, 0x33) # Dark Grey text
    TITLE_COLOR = RGBColor(0x00, 0x00, 0x00) # Black for titles
    ACCENT_COLOR = RGBColor(0x6A, 0x0D, 0xAD) # A professional purple
    FOOTER_TEXT = "GetAligned | Confidential"

    # --- Slide 0: Main Title Slide ---
    slide_layout = prs.slide_layouts[5] # Using a title-only layout as a base
    slide = prs.slides.add_slide(slide_layout)
    
    # Set gradient background for title slide
    set_slide_background_gradient(slide, RGBColor(0xF8, 0xF9, 0xFA), RGBColor(0xE9, 0xEF, 0xF7))
    
    # Add decorative shapes
    add_decorative_shapes(slide)
    
    # Add title background banner
    add_title_background_banner(slide, data["title"])
    
    # Add icon for the title slide
    add_icon_shape(slide, "lightbulb", 1, 2, 0.8)

    # Main Title
    title = slide.shapes.title
    title.text = data["title"]
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    core_problem = data['strategic_plan']['customer_analysis']['core_problem']
    customer_name = core_problem.split("'s")[0]
    p.text = f"A Strategic Presentation for {customer_name}"
    p.font.color.rgb = TEXT_COLOR
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    add_footer(slide, FOOTER_TEXT)

    # --- Generate Slides from JSON Data ---
    slide_icons = ["gear", "lightbulb", "star", "check", "gear", "star", "check", "gear", "lightbulb", "star", "check", "gear"]
    
    for i, slide_data in enumerate(data["slides"]):
        slide_layout = prs.slide_layouts[5] # Title Only layout as a base
        slide = prs.slides.add_slide(slide_layout)
        
        # Alternate between gradient and solid backgrounds
        if i % 2 == 0:
            set_slide_background_gradient(slide, RGBColor(0xF8, 0xF9, 0xFA), RGBColor(0xE9, 0xEF, 0xF7))
        else:
            set_slide_background(slide, BG_COLOR)
        
        # Add decorative shapes
        add_decorative_shapes(slide)
        
        # Add title background banner
        add_title_background_banner(slide, slide_data["title"])
        
        # Add relevant icon for each slide
        icon_type = slide_icons[i % len(slide_icons)]
        add_icon_shape(slide, icon_type, 0.5, 1.8, 0.6)
        
        # Slide Title
        title = slide.shapes.title
        title.text = slide_data["title"]
        title.top = Inches(0.5)
        title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Slide Content (Key Elements)
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7.5), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for j, element in enumerate(slide_data["key_content_elements"]):
            p = tf.add_paragraph()
            p.text = f"• {element}"  # Add bullet points
            p.font.color.rgb = TEXT_COLOR
            p.font.size = Pt(18)
            p.space_before = Pt(12)
            p.level = 0
            
            # Add small icons next to key points
            if j < 3:  # Only for first 3 points to avoid clutter
                add_icon_shape(slide, "check", 1.2, 2.2 + (j * 0.8), 0.3)

        add_footer(slide, FOOTER_TEXT)

        # Add presenter notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"Objective: {slide_data['objective']}\n\nTalking Points:\n{slide_data['sales_narrative']}"

    # --- Final Slide: Call to Action ---
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    
    # Special gradient for CTA slide
    set_slide_background_gradient(slide, RGBColor(0x6A, 0x0D, 0xAD), RGBColor(0x9B, 0x59, 0xB6))
    
    # Add decorative shapes
    add_decorative_shapes(slide)
    
    # Add CTA icon
    add_icon_shape(slide, "star", 4.5, 2.5, 1.0)

    # Call to Action Title
    title = slide.shapes.title
    title.text = "Next Steps"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White for contrast
    title.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(36)

    # Call to Action Text Box with enhanced styling
    cta_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(3.5),
        Inches(8), Inches(2)
    )
    cta_bg.fill.solid()
    cta_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cta_bg.fill.transparency = 0.1
    cta_bg.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cta_bg.line.width = Pt(3)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data['strategic_plan']['primary_cta']
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_footer(slide, FOOTER_TEXT)

    # Save the presentation
    file_path = "GetAligned_Sales_Presentation.pptx"
    prs.save(file_path)
    return file_path

# --- Main Execution ---
# The JSON data provided by the user.
input_data = {
    "id": 67890,
    "slides": [
        {
           "id": 1, "title": "The Hidden Cost of Manual Sales Processes", "objective": "Highlight pain points", "key_content_elements": ["Headline focusing on time waste and inefficiency.", "Data points on time spent on manual tasks (4-6 hours/week).", "Mention of inconsistent documentation and lost insights."], "sales_narrative": "Priya, Aaron, Ravi, we understand the pressures... leading to decreased productivity and missed opportunities.", "visual_concept": "A visually compelling image or graphic representing wasted time and resources."
       }, {
           "id": 2, "title": "Introducing GetAligned: Your AI Co-Pilot", "objective": "Position GetAligned as the solution.", "key_content_elements": ["Headline introducing GetAligned as the solution.", "A brief overview of GetAligned's key features.", "A concise tagline that captures GetAligned's value proposition."], "sales_narrative": "GetAligned is an AI-driven software designed to automate insight extraction... unlock the full potential of your sales conversations.", "visual_concept": "Clean and modern design with a clear visual representation of GetAligned's functionality."
       }, {
           "id": 3, "title": "Automated Insights, Streamlined Workflows", "objective": "Showcase specific features and benefits.", "key_content_elements": ["Real-time summary generation.", "Automated pitch deck creation.", "Quantifiable benefits (time savings, increased productivity).", "Integration capabilities (HubSpot, Slack, Google Drive)."], "sales_narrative": "Remember those 4-6 hours reps spend on manual updates? GetAligned automates that... imagine the impact on your RevOps team.", "visual_concept": "Screenshots and visuals demonstrating GetAligned's features."
       }, {
           "id": 4, "title": "Addressing Your Concerns", "objective": "Address objections head-on.", "key_content_elements": ["Clear answers to objections (data privacy, accuracy, change management).", "Info on security certifications.", "Details on onboarding and support."], "sales_narrative": "Priya, you raised valid concerns about data privacy... ensure a smooth transition for your team.", "visual_concept": "Icons representing security, accuracy, and support."
       }, {
           "id": 5, "title": "Powering Success: Real-World Results", "objective": "Provide social proof.", "key_content_elements": ["Customer testimonials highlighting benefits.", "Case studies demonstrating impact.", "Quantifiable results (increased sales, reduced costs).", "Logos of satisfied customers."], "sales_narrative": "Don't just take our word for it. Companies like [Similar Company] have seen a [Result] increase in sales...", "visual_concept": "Customer logos and quotes displayed prominently."
       }, {
           "id": 6, "title": "Pilot Program: Experience the Power Risk-Free", "objective": "Present the pilot offer.", "key_content_elements": ["Details of the 2-month free pilot for up to 5 users.", "Clear explanation of the pilot benefits.", "A call to action to sign up."], "sales_narrative": "To demonstrate value, we're offering a 2-month free pilot... see how it can transform your sales process.", "visual_concept": "A graphic highlighting the value of the pilot program."
       }, {
           "id": 7, "title": "Pricing & ROI: Invest in Sales Excellence", "objective": "Present pricing and demonstrate ROI.", "key_content_elements": ["Seat-based, monthly subscription model.", "Estimated cost: $49/user/month.", "Available discounts (annual, enterprise).", "ROI calculations showing savings and revenue gains."], "sales_narrative": "Our pricing is flexible and scalable... Even 2 hours per rep is worth piloting.", "visual_concept": "A table comparing the cost of GetAligned to manual processes."
       }, {
           "id": 8, "title": "Seamless Integration", "objective": "Showcase integration capabilities.", "key_content_elements": ["Logos of integration partners (HubSpot, Slack, Google Drive).", "Brief explanation of the integration process.", "Highlight ease of integration."], "sales_narrative": "GetAligned seamlessly integrates with your existing tools... without disrupting your current workflow.", "visual_concept": "A diagram illustrating integrations."
       }, {
           "id": 9, "title": "Onboarding & Support", "objective": "Describe the support process.", "key_content_elements": ["Live onboarding session for team leads.", "Available support channels (Slack, email, chat).", "Response time SLA (within 4 business hours).", "Training materials (videos, playbooks)."], "sales_narrative": "We're committed to your success... get everyone up to speed quickly.", "visual_concept": "A graphic showcasing the onboarding process."
       }, {
           "id": 10, "title": "Next Steps: Let's Move Forward", "objective": "Outline next steps clearly.", "key_content_elements": ["List of next steps (technical deep-dive, share sample summary).", "Call to action to schedule the deep-dive.", "Timeline for next steps."], "sales_narrative": "The next step is to schedule a technical deep-dive... Are you available next week for a follow up?", "visual_concept": "A timeline illustrating the next steps."
       }, {
           "id": 11, "title": "Thank You & Q&A", "objective": "Open the floor for questions.", "key_content_elements": ["A thank you message.", "A clear invitation to ask questions.", "Contact information for follow-up."], "sales_narrative": "Thank you for your time... What questions do you have?", "visual_concept": "A simple and professional slide with contact information."
       }, {
           "id": 12, "title": "Appendix: Competitive Analysis", "objective": "Address the competitive landscape.", "key_content_elements": ["Comparison table to competitors.", "Highlight GetAligned's unique advantages.", "Key differentiators (Automated Pitch Deck, HubSpot Integration)."], "sales_narrative": "We understand you're evaluating other solutions... maximizes the value of your sales conversations.", "visual_concept": "A comparison table."
       }
    ],
    "strategic_plan": {
       "customer_analysis": {
           "core_problem": "FinCore's sales team is losing significant time..."
       },
       "primary_cta": "Schedule a technical deep-dive with the RevOps team."
    },
    "title": "Unlock Sales Excellence: AI-Powered Insights for FinCore Solutions"
}


# Generate the presentation
generated_file = create_presentation(input_data)
print(f"Presentation has been generated: {generated_file}")
