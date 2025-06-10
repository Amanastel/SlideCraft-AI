#!/usr/bin/env python3
"""
Advanced PowerPoint Generator with LLM Integration and Image Support
This script generates presentations based on prompts and structured data,
with automatic image integration from various sources.
"""

import json
import requests
from io import BytesIO
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os
from typing import Dict, List, Optional
import time

# Try to import OpenAI - make it optional
try:
    import openai
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False
    print("⚠️  OpenAI package not installed. LLM features will be disabled.")

class AdvancedPPTGenerator:
    def __init__(self, openai_api_key: Optional[str] = None):
        """Initialize the PPT generator with optional OpenAI API key."""
        self.openai_api_key = openai_api_key or os.getenv('OPENAI_API_KEY')
        if self.openai_api_key and OPENAI_AVAILABLE:
            openai.api_key = self.openai_api_key
        elif self.openai_api_key and not OPENAI_AVAILABLE:
            print("⚠️  OpenAI API key provided but OpenAI package not installed.")
        
        # Image sources and APIs
        self.unsplash_access_key = os.getenv('UNSPLASH_ACCESS_KEY')
        self.pixabay_api_key = os.getenv('PIXABAY_API_KEY')
        
        # Color schemes
        self.color_schemes = {
            'professional': {
                'primary': RGBColor(0x6A, 0x0D, 0xAD),
                'secondary': RGBColor(0x9B, 0x59, 0xB6),
                'accent': RGBColor(0x3D, 0xBD, 0xE0),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF1, 0xC4, 0x0F),
                'text_dark': RGBColor(0x33, 0x33, 0x33),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF),
                'background': RGBColor(0xFA, 0xFA, 0xFA)
            }
        }

    def enhance_content_with_llm(self, prompt: str, slide_data: Dict) -> Dict:
        """Use LLM to enhance slide content based on prompt."""
        if not self.openai_api_key or not OPENAI_AVAILABLE:
            print("⚠️  No OpenAI API key provided or OpenAI not available. Using original content.")
            return slide_data
        
        try:
            enhancement_prompt = f"""
            Based on this presentation data and user prompt, enhance the slide content:
            
            User Prompt: {prompt}
            
            Original Slide Data: {json.dumps(slide_data, indent=2)}
            
            Please enhance the content by:
            1. Making titles more engaging and specific
            2. Improving key content elements to be more compelling
            3. Suggesting relevant visual concepts
            4. Adding suggested image search terms for each slide
            5. Maintaining the original structure and intent
            
            Return the enhanced data in the same JSON format with additional fields:
            - enhanced_title (if different)
            - enhanced_key_content_elements
            - image_search_terms (array of search terms for images)
            - visual_style_suggestions
            """
            
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a presentation design expert who creates compelling slide content."},
                    {"role": "user", "content": enhancement_prompt}
                ],
                max_tokens=1500,
                temperature=0.7
            )
            
            enhanced_content = response.choices[0].message.content
            # Try to parse JSON response
            try:
                enhanced_data = json.loads(enhanced_content)
                return {**slide_data, **enhanced_data}
            except json.JSONDecodeError:
                print(f"⚠️  Could not parse LLM response as JSON. Using original content.")
                return slide_data
                
        except Exception as e:
            print(f"⚠️  LLM enhancement failed: {e}")
            return slide_data

    def search_unsplash_image(self, query: str, per_page: int = 3) -> List[str]:
        """Search for images on Unsplash."""
        if not self.unsplash_access_key:
            return []
        
        try:
            url = "https://api.unsplash.com/search/photos"
            headers = {"Authorization": f"Client-ID {self.unsplash_access_key}"}
            params = {
                "query": query,
                "per_page": per_page,
                "orientation": "landscape"
            }
            
            response = requests.get(url, headers=headers, params=params, timeout=10)
            if response.status_code == 200:
                data = response.json()
                return [photo["urls"]["regular"] for photo in data.get("results", [])]
        except Exception as e:
            print(f"⚠️  Unsplash search failed for '{query}': {e}")
        
        return []

    def search_pixabay_image(self, query: str, per_page: int = 3) -> List[str]:
        """Search for images on Pixabay."""
        if not self.pixabay_api_key:
            return []
        
        try:
            url = "https://pixabay.com/api/"
            params = {
                "key": self.pixabay_api_key,
                "q": query,
                "image_type": "photo",
                "orientation": "horizontal",
                "per_page": per_page,
                "safesearch": "true"
            }
            
            response = requests.get(url, params=params, timeout=10)
            if response.status_code == 200:
                data = response.json()
                return [img["webformatURL"] for img in data.get("hits", [])]
        except Exception as e:
            print(f"⚠️  Pixabay search failed for '{query}': {e}")
        
        return []

    def get_fallback_images(self, category: str) -> List[str]:
        """Get fallback images from public APIs that don't require keys."""
        fallback_urls = {
            'business': [
                'https://picsum.photos/1200/600?random=1',
                'https://picsum.photos/1200/600?random=2'
            ],
            'technology': [
                'https://picsum.photos/1200/600?random=3',
                'https://picsum.photos/1200/600?random=4'
            ],
            'success': [
                'https://picsum.photos/1200/600?random=5',
                'https://picsum.photos/1200/600?random=6'
            ],
            'abstract': [
                'https://picsum.photos/1200/600?random=7',
                'https://picsum.photos/1200/600?random=8'
            ]
        }
        return fallback_urls.get(category, fallback_urls['abstract'])

    def download_image(self, url: str) -> Optional[BytesIO]:
        """Download image from URL and return as BytesIO object."""
        try:
            response = requests.get(url, timeout=15)
            if response.status_code == 200:
                return BytesIO(response.content)
        except Exception as e:
            print(f"⚠️  Failed to download image from {url}: {e}")
        return None

    def add_background_image(self, slide, image_url: str, opacity: float = 0.3):
        """Add background image to slide with opacity."""
        image_stream = self.download_image(image_url)
        if image_stream:
            try:
                # Add image to slide
                pic = slide.shapes.add_picture(
                    image_stream, 
                    Inches(0), Inches(0),
                    Inches(10), Inches(7.5)
                )
                
                # Move to background
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
                
                # Apply transparency
                pic.fill.transparency = opacity
                
                return True
            except Exception as e:
                print(f"⚠️  Failed to add background image: {e}")
        return False

    def add_content_image(self, slide, image_url: str, x: float, y: float, width: float, height: float):
        """Add content image to slide at specified position."""
        image_stream = self.download_image(image_url)
        if image_stream:
            try:
                pic = slide.shapes.add_picture(
                    image_stream,
                    Inches(x), Inches(y),
                    Inches(width), Inches(height)
                )
                return pic
            except Exception as e:
                print(f"⚠️  Failed to add content image: {e}")
        return None

    def create_enhanced_slide_background(self, slide, style: str = 'gradient'):
        """Create enhanced slide backgrounds."""
        colors = self.color_schemes['professional']
        
        if style == 'gradient':
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = colors['background']
            fill.gradient_stops[1].color.rgb = RGBColor(0xE9, 0xEF, 0xF7)
        elif style == 'accent_gradient':
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = colors['primary']
            fill.gradient_stops[1].color.rgb = colors['secondary']
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = colors['background']

    def add_decorative_elements(self, slide, theme: str = 'professional'):
        """Add decorative elements to slide."""
        colors = self.color_schemes['professional']
        
        # Top right accent
        shape1 = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(8.5), Inches(0), 
            Inches(1.5), Inches(1.5)
        )
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = colors['primary']
        shape1.fill.transparency = 0.8
        shape1.line.fill.background()
        
        # Bottom left accent
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            Inches(-0.5), Inches(6.5), 
            Inches(2), Inches(2)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = colors['accent']
        shape2.fill.transparency = 0.9
        shape2.line.fill.background()

    def add_data_chart(self, slide, chart_data: Dict, chart_type: str = 'bar'):
        """Add data visualization chart to slide."""
        try:
            # Create chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get('categories', ['Category 1', 'Category 2', 'Category 3'])
            chart_data_obj.add_series('Data', chart_data.get('values', [10, 20, 30]))
            
            # Determine chart type
            chart_type_map = {
                'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE
            }
            
            xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Add chart to slide
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
            chart = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, chart_data_obj).chart
            
            return chart
        except Exception as e:
            print(f"⚠️  Failed to add chart: {e}")
            return None

    def create_enhanced_footer(self, slide, text: str, colors: Dict = None):
        """Create enhanced footer with dynamic colors."""
        if colors is None:
            colors = self.color_schemes['professional']
            
        # Footer background with dynamic color
        footer_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.8),
            Inches(10), Inches(0.7)
        )
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = colors['text_dark']
        footer_bg.line.fill.background()
        
        # Footer text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.color.rgb = colors['text_light']
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # Dynamic logo placeholder
        logo_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5), Inches(7.0),
            Inches(1), Inches(0.4)
        )
        logo_shape.fill.solid()
        logo_shape.fill.fore_color.rgb = colors['primary']
        logo_shape.line.fill.background()

    def generate_presentation_from_prompt(self, prompt: str, input_data: Dict, output_filename: str = None) -> str:
        """Generate enhanced presentation from prompt and flexible data structure."""
        print("🚀 Starting Advanced PPT Generation...")
        print(f"📝 User Prompt: {prompt}")
        
        # Normalize input data to handle flexible structures
        print("🔄 Normalizing input data structure...")
        normalized_data = self.normalize_presentation_data(input_data)
        print(f"📊 Normalized {len(normalized_data.get('slides', []))} slides")
        
        # Initialize presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        colors = self.color_schemes['professional']
        
        # Generate title slide
        print("📄 Creating title slide...")
        self._create_title_slide(prs, normalized_data, prompt)
        
        # Generate content slides
        slides_data = normalized_data.get('slides', [])
        total_slides = len(slides_data)
        
        for i, slide_data in enumerate(slides_data, 1):
            print(f"📄 Creating slide {i}/{total_slides}: {slide_data.get('title', 'Untitled')}")
            
            # Enhance content with LLM
            enhanced_slide_data = self.enhance_content_with_llm(prompt, slide_data)
            
            # Create slide
            self._create_content_slide(prs, enhanced_slide_data, i)
            
            # Add brief delay to avoid rate limiting
            time.sleep(0.1)
        
        # Add call-to-action slide
        print("📄 Creating call-to-action slide...")
        self._create_cta_slide(prs, normalized_data)
        
        # Save presentation
        output_file = output_filename or "Advanced_AI_Generated_Presentation.pptx"
        prs.save(output_file)
        
        print(f"✅ Presentation saved as: {output_file}")
        return output_file

    def analyze_slide_content(self, slide_data: Dict) -> Dict:
        """
        Analyze slide content to determine appropriate styling, icons, and layout.
        Returns content analysis with recommendations.
        """
        title = slide_data.get('title', '').lower()
        content_elements = slide_data.get('key_content_elements', [])
        content_text = ' '.join([str(elem) for elem in content_elements]).lower()
        all_text = f"{title} {content_text}".lower()
        
        analysis = {
            'content_type': 'general',
            'emotional_tone': 'neutral',
            'data_heavy': False,
            'action_oriented': False,
            'visual_focus': 'balanced',
            'recommended_layout': 'standard',
            'color_scheme': 'professional',
            'image_size': 'medium',
            'background_style': 'gradient'
        }
        
        # Analyze content type
        if any(word in all_text for word in ['problem', 'challenge', 'issue', 'difficulty', 'pain']):
            analysis['content_type'] = 'problem'
            analysis['emotional_tone'] = 'concerned'
            analysis['color_scheme'] = 'attention'
            
        elif any(word in all_text for word in ['solution', 'fix', 'resolve', 'answer', 'innovation']):
            analysis['content_type'] = 'solution'
            analysis['emotional_tone'] = 'positive'
            analysis['color_scheme'] = 'success'
            
        elif any(word in all_text for word in ['result', 'outcome', 'achievement', 'success', 'roi', 'benefit']):
            analysis['content_type'] = 'results'
            analysis['emotional_tone'] = 'triumphant'
            analysis['color_scheme'] = 'achievement'
            
        elif any(word in all_text for word in ['feature', 'capability', 'function', 'specification']):
            analysis['content_type'] = 'features'
            analysis['color_scheme'] = 'technical'
            
        elif any(word in all_text for word in ['process', 'step', 'workflow', 'methodology']):
            analysis['content_type'] = 'process'
            analysis['recommended_layout'] = 'timeline'
            
        elif any(word in all_text for word in ['team', 'people', 'user', 'customer', 'client']):
            analysis['content_type'] = 'people'
            analysis['color_scheme'] = 'warm'
            
        elif any(word in all_text for word in ['data', 'analytics', 'metrics', 'statistics', 'numbers']):
            analysis['content_type'] = 'data'
            analysis['data_heavy'] = True
            analysis['recommended_layout'] = 'data_focused'
            analysis['visual_focus'] = 'charts'
            
        # Analyze if action-oriented
        if any(word in all_text for word in ['action', 'next steps', 'implement', 'execute', 'start', 'begin']):
            analysis['action_oriented'] = True
            analysis['background_style'] = 'dynamic'
            
        # Determine image sizing based on content
        if analysis['data_heavy']:
            analysis['image_size'] = 'small'  # Leave more space for charts
        elif analysis['visual_focus'] == 'charts':
            analysis['image_size'] = 'small'
        elif analysis['content_type'] in ['problem', 'solution']:
            analysis['image_size'] = 'large'  # Visual impact
        elif len(content_elements) > 5:
            analysis['image_size'] = 'medium'  # Balanced with text
        else:
            analysis['image_size'] = 'large'
            
        return analysis

    def get_dynamic_color_scheme(self, scheme_name: str) -> Dict:
        """Get dynamic color schemes based on content analysis."""
        
        schemes = {
            'professional': {
                'primary': RGBColor(0x6A, 0x0D, 0xAD),
                'secondary': RGBColor(0x9B, 0x59, 0xB6),
                'accent': RGBColor(0x3D, 0xBD, 0xE0),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xF8, 0xF9, 0xFA),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            },
            'attention': {
                'primary': RGBColor(0xE7, 0x4C, 0x3C),
                'secondary': RGBColor(0xF3, 0x9C, 0x12),
                'accent': RGBColor(0xFF, 0x6B, 0x6B),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xFD, 0xF2, 0xF2),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            },
            'success': {
                'primary': RGBColor(0x2E, 0xCC, 0x71),
                'secondary': RGBColor(0x1D, 0xD1, 0xA1),
                'accent': RGBColor(0x55, 0xEF, 0xC4),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xF0, 0xFD, 0xF4),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            },
            'achievement': {
                'primary': RGBColor(0xF3, 0x9C, 0x12),
                'secondary': RGBColor(0xE6, 0x7E, 0x22),
                'accent': RGBColor(0xFF, 0xD7, 0x00),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xFF, 0xF9, 0xE6),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            },
            'technical': {
                'primary': RGBColor(0x3D, 0xBD, 0xE0),
                'secondary': RGBColor(0x5D, 0xAD, 0xE2),
                'accent': RGBColor(0x85, 0xC1, 0xE9),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xF4, 0xF8, 0xFB),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            },
            'warm': {
                'primary': RGBColor(0xFF, 0x7F, 0x50),
                'secondary': RGBColor(0xFF, 0x6B, 0x6B),
                'accent': RGBColor(0xFF, 0xA0, 0x7A),
                'success': RGBColor(0x2E, 0xCC, 0x71),
                'warning': RGBColor(0xF3, 0x9C, 0x12),
                'danger': RGBColor(0xE7, 0x4C, 0x3C),
                'background': RGBColor(0xFF, 0xF5, 0xF5),
                'text_dark': RGBColor(0x2C, 0x3E, 0x50),
                'text_light': RGBColor(0xFF, 0xFF, 0xFF)
            }
        }
        
        return schemes.get(scheme_name, schemes['professional'])

    def get_image_dimensions(self, image_size: str, slide_layout: str = 'standard') -> Dict:
        """Get dynamic image dimensions based on content analysis."""
        
        dimensions = {
            'small': {
                'width': 2.0,
                'height': 1.5,
                'x': 7.5,
                'y': 2.0
            },
            'medium': {
                'width': 3.0,
                'height': 2.25,
                'x': 6.5,
                'y': 1.5
            },
            'large': {
                'width': 4.0,
                'height': 3.0,
                'x': 5.5,
                'y': 1.0
            }
        }
        
        # Adjust for different layouts
        if slide_layout == 'data_focused':
            # Smaller images for data-heavy slides
            for size in dimensions:
                dimensions[size]['width'] *= 0.7
                dimensions[size]['height'] *= 0.7
                dimensions[size]['x'] += 0.5
                
        elif slide_layout == 'timeline':
            # Horizontal adjustment for process slides
            for size in dimensions:
                dimensions[size]['width'] *= 1.2
                dimensions[size]['height'] *= 0.8
                dimensions[size]['y'] += 0.5
                
        return dimensions.get(image_size, dimensions['medium'])

    def create_dynamic_background(self, slide, analysis: Dict, colors: Dict):
        """Create dynamic backgrounds based on content analysis."""
        
        background_style = analysis.get('background_style', 'gradient')
        content_type = analysis.get('content_type', 'general')
        
        background = slide.background
        fill = background.fill
        
        if background_style == 'dynamic' and analysis.get('action_oriented'):
            # Diagonal gradient for action slides
            fill.gradient()
            fill.gradient_angle = 45
            fill.gradient_stops[0].color.rgb = colors['primary']
            fill.gradient_stops[1].color.rgb = colors['accent']
            
        elif content_type == 'problem':
            # Subtle warning background
            fill.gradient()
            fill.gradient_stops[0].color.rgb = colors['background']
            fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
            
        elif content_type == 'success' or content_type == 'results':
            # Success gradient
            fill.gradient()
            fill.gradient_stops[0].color.rgb = colors['background']
            fill.gradient_stops[1].color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
            
        elif content_type == 'data':
            # Clean, minimal background for data
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)
            
        else:
            # Default professional gradient
            fill.gradient()
            fill.gradient_stops[0].color.rgb = colors['background']
            fill.gradient_stops[1].color.rgb = RGBColor(0xE9, 0xEF, 0xF7)

    def add_smart_decorative_elements(self, slide, analysis: Dict, colors: Dict):
        """Add contextually appropriate decorative elements."""
        
        content_type = analysis.get('content_type', 'general')
        emotional_tone = analysis.get('emotional_tone', 'neutral')
        
        if content_type == 'problem':
            # Subtle warning accent
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, 
                Inches(8.5), Inches(0.2), 
                Inches(1.0), Inches(1.0)
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = colors['warning']
            shape1.fill.transparency = 0.7
            shape1.line.fill.background()
            
        elif content_type == 'solution':
            # Innovation burst
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.STAR_4_POINT, 
                Inches(8.7), Inches(0.2), 
                Inches(0.8), Inches(0.8)
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = colors['success']
            shape1.fill.transparency = 0.6
            shape1.line.fill.background()
            
        elif content_type == 'results' or content_type == 'success':
            # Achievement crown
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.EXPLOSION1, 
                Inches(8.6), Inches(0.1), 
                Inches(1.0), Inches(1.0)
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = colors['warning']
            shape1.fill.transparency = 0.5
            shape1.line.fill.background()
            
        elif content_type == 'data':
            # Clean geometric accent
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON, 
                Inches(8.7), Inches(0.3), 
                Inches(0.6), Inches(0.6)
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = colors['accent']
            shape1.fill.transparency = 0.8
            shape1.line.fill.background()
            
        else:
            # Default professional accent
            shape1 = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(8.5), Inches(0), 
                Inches(1.5), Inches(1.2)
            )
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = colors['primary']
            shape1.fill.transparency = 0.8
            shape1.line.fill.background()
        
        # Bottom accent based on action orientation
        if analysis.get('action_oriented'):
            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, 
                Inches(-0.3), Inches(6.8), 
                Inches(1.5), Inches(0.7)
            )
            shape2.fill.solid()
            shape2.fill.fore_color.rgb = colors['accent']
            shape2.fill.transparency = 0.9
            shape2.line.fill.background()
        else:
            shape2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, 
                Inches(-0.5), Inches(6.5), 
                Inches(2), Inches(2)
            )
            shape2.fill.solid()
            shape2.fill.fore_color.rgb = colors['accent']
            shape2.fill.transparency = 0.9
            shape2.line.fill.background()

    def add_adaptive_chart(self, slide, chart_data: Dict, x: float, width: float, colors: Dict):
        """Add chart with adaptive positioning and styling."""
        try:
            # Create chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get('categories', ['Category 1', 'Category 2', 'Category 3'])
            chart_data_obj.add_series('Data', chart_data.get('values', [10, 20, 30]))
            
            # Smart chart type selection
            values = chart_data.get('values', [10, 20, 30])
            if len(values) <= 5 and all(v > 0 for v in values):
                chart_type = XL_CHART_TYPE.PIE  # Pie for small positive datasets
            elif any('trend' in str(cat).lower() or 'time' in str(cat).lower() 
                    for cat in chart_data.get('categories', [])):
                chart_type = XL_CHART_TYPE.LINE  # Line for trends
            else:
                chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED  # Bar for comparisons
            
            # Adaptive positioning
            chart_y = 4.0 if len(chart_data.get('categories', [])) > 4 else 3.5
            chart_height = 2.5 if len(chart_data.get('categories', [])) > 4 else 3.0
            
            # Add chart to slide
            chart_shape = slide.shapes.add_chart(
                chart_type, 
                Inches(x), Inches(chart_y), 
                Inches(width), Inches(chart_height), 
                chart_data_obj
            )
            
            chart = chart_shape.chart
            
            # Style the chart with dynamic colors
            if hasattr(chart, 'series'):
                series = chart.series[0]
                if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = colors['primary']
            
            return chart
        except Exception as e:
            print(f"⚠️  Failed to add adaptive chart: {e}")
            return None

    def normalize_slide_data(self, slide_data: Dict) -> Dict:
        """
        Convert any slide structure to the expected format.
        Handles flexible JSON input structures like the user's example.
        """
        normalized = {}
        
        # Handle title
        normalized['title'] = slide_data.get('title', 'Untitled Slide')
        
        # Handle content - can be 'content', 'key_content_elements', or other formats
        content_elements = []
        
        if 'content' in slide_data:
            # Split content by newlines and bullet points
            content = slide_data['content']
            if isinstance(content, str):
                # Split by newlines and clean up
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                for line in lines:
                    # Remove bullet points and dashes
                    clean_line = line.lstrip('•-* ').strip()
                    if clean_line:
                        content_elements.append(clean_line)
            elif isinstance(content, list):
                content_elements = content
        
        if 'key_content_elements' in slide_data:
            # Merge with existing content elements
            if isinstance(slide_data['key_content_elements'], list):
                content_elements.extend(slide_data['key_content_elements'])
        
        # If no content elements found, create from title
        if not content_elements:
            content_elements = [f"Key information about {normalized['title']}"]
        
        normalized['key_content_elements'] = content_elements[:6]  # Limit to 6 elements
        
        # Handle image references
        image_search_terms = []
        
        if 'image' in slide_data:
            # Extract keywords from image filename/path
            image_name = slide_data['image']
            if isinstance(image_name, str):
                # Remove file extensions and split by common separators
                image_keywords = image_name.replace('.jpg', '').replace('.png', '').replace('.jpeg', '')
                image_keywords = image_keywords.replace('-', ' ').replace('_', ' ')
                image_search_terms.append(image_keywords)
        
        if 'image_search_terms' in slide_data:
            if isinstance(slide_data['image_search_terms'], list):
                image_search_terms.extend(slide_data['image_search_terms'])
            elif isinstance(slide_data['image_search_terms'], str):
                image_search_terms.append(slide_data['image_search_terms'])
        
        # Generate search terms from title and content if none provided
        if not image_search_terms:
            title_words = normalized['title'].lower().split()
            content_words = ' '.join(content_elements[:2]).lower().split()
            
            # Extract meaningful keywords (longer than 3 characters)
            keywords = [word for word in title_words + content_words 
                       if len(word) > 3 and word not in ['the', 'and', 'for', 'with']]
            
            if keywords:
                image_search_terms = [' '.join(keywords[:3])]  # Take first 3 keywords
            else:
                image_search_terms = ['business professional']
        
        normalized['image_search_terms'] = image_search_terms
        
        # Handle chart data if present
        if 'chart_data' in slide_data:
            normalized['chart_data'] = slide_data['chart_data']
        
        # Set objective based on content
        if 'objective' in slide_data:
            normalized['objective'] = slide_data['objective']
        else:
            # Generate objective from title
            normalized['objective'] = f"Present information about {normalized['title'].lower()}"
        
        # Handle any additional custom fields
        for key, value in slide_data.items():
            if key not in ['title', 'content', 'key_content_elements', 'image', 'image_search_terms', 'chart_data', 'objective']:
                normalized[f'custom_{key}'] = value
        
        return normalized

    def normalize_presentation_data(self, input_data: Dict) -> Dict:
        """
        Convert any presentation structure to the expected format.
        Handles flexible JSON input structures.
        """
        normalized = {}
        
        # Handle title
        normalized['title'] = input_data.get('title', 'Presentation')
        
        # Handle slides
        slides = input_data.get('slides', [])
        normalized_slides = []
        
        for i, slide_data in enumerate(slides):
            normalized_slide = self.normalize_slide_data(slide_data)
            normalized_slide['id'] = i + 1  # Ensure slides have IDs
            normalized_slides.append(normalized_slide)
        
        normalized['slides'] = normalized_slides
        
        # Handle strategic plan / call-to-action
        strategic_plan = {}
        
        if 'strategic_plan' in input_data:
            strategic_plan = input_data['strategic_plan']
        else:
            # Create default strategic plan
            strategic_plan = {
                'primary_cta': 'Contact us to learn more about our solution',
                'target_audience': 'Business professionals and decision makers',
                'goal': 'Present compelling information and drive action'
            }
        
        normalized['strategic_plan'] = strategic_plan
        
        # Handle any additional metadata
        for key, value in input_data.items():
            if key not in ['title', 'slides', 'strategic_plan']:
                normalized[f'metadata_{key}'] = value
        
        return normalized

    def _create_title_slide(self, prs, input_data: Dict, prompt: str):
        """Create enhanced title slide with dynamic styling."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Analyze overall presentation for title slide styling
        title = input_data.get('title', 'Presentation Title')
        title_analysis = self.analyze_slide_content({'title': title, 'key_content_elements': [prompt]})
        colors = self.get_dynamic_color_scheme(title_analysis['color_scheme'])
        
        # Create dramatic background for title slide
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135  # Diagonal gradient
        fill.gradient_stops[0].color.rgb = colors['primary']
        fill.gradient_stops[1].color.rgb = colors['secondary']
        
        # Search for hero background image
        bg_search_terms = f"{title} professional presentation business"
        bg_images = self.search_unsplash_image(bg_search_terms) or self.get_fallback_images('business')
        if bg_images:
            self.add_background_image(slide, bg_images[0], opacity=0.15)
        
        # Title with adaptive sizing
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_shape.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
        
        # Dynamic title size based on length
        title_length = len(title)
        if title_length < 20:
            title_size = 48
        elif title_length < 40:
            title_size = 40
        else:
            title_size = 32
            
        title_shape.text_frame.paragraphs[0].font.size = Pt(title_size)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Enhanced subtitle with dynamic content
        subtitle_content = f"AI-Enhanced Presentation | {title_analysis['content_type'].title()} Focus"
        if len(prompt) > 50:
            subtitle_content += f"\nPrompt: {prompt[:80]}..."
        else:
            subtitle_content += f"\nPrompt: {prompt}"
            
        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle_content
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT
        p.font.bold = False
        
        # Dynamic footer for title slide
        footer_text = f"Generated with Advanced AI | {title_analysis['emotional_tone'].title()} Tone"
        self.create_enhanced_footer(slide, footer_text)

    def _create_content_slide(self, prs, slide_data: Dict, slide_number: int):
        """Create enhanced content slide with dynamic styling based on content analysis."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Analyze slide content for smart styling
        content_analysis = self.analyze_slide_content(slide_data)
        print(f"   📊 Content analysis: {content_analysis['content_type']} slide, dynamic styling applied")
        
        # Get dynamic color scheme
        colors = self.get_dynamic_color_scheme(content_analysis['color_scheme'])
        
        # Create dynamic background
        self.create_dynamic_background(slide, content_analysis, colors)
        
        # Search for relevant images with smart sizing
        title = slide_data.get('title', '')
        search_terms = slide_data.get('image_search_terms', [title.split()[:3]])
        image_dims = self.get_image_dimensions(
            content_analysis['image_size'], 
            content_analysis['recommended_layout']
        )
        
        if search_terms:
            search_query = ' '.join(search_terms[0]) if isinstance(search_terms[0], list) else search_terms[0]
            content_images = self.search_unsplash_image(search_query) or self.get_fallback_images('business')
            if content_images:
                self.add_content_image(
                    slide, 
                    content_images[0], 
                    image_dims['x'], 
                    image_dims['y'],
                    image_dims['width'], 
                    image_dims['height']
                )
        
        # Add smart decorative elements based on content
        self.add_smart_decorative_elements(slide, content_analysis, colors)
        
        # Title with dynamic positioning
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get('enhanced_title', slide_data.get('title', 'Slide Title'))
        title_shape.top = Inches(0.5)
        title_shape.text_frame.paragraphs[0].font.color.rgb = colors['text_dark']
        title_shape.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Content with adaptive layout
        content_elements = slide_data.get('enhanced_key_content_elements', slide_data.get('key_content_elements', []))
        
        # Handle different content structures
        if not content_elements and 'content' in slide_data:
            content = slide_data['content']
            if isinstance(content, str):
                content_elements = [line.strip().lstrip('•-* ') for line in content.split('\n') if line.strip()]
            elif isinstance(content, list):
                content_elements = content
        
        # Adjust content box based on image size and layout
        if content_analysis['image_size'] == 'large':
            content_width = 4.0
        elif content_analysis['image_size'] == 'medium':
            content_width = 5.0
        else:
            content_width = 6.0
            
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(content_width), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Add content elements with smart icons
        max_elements = 6 if content_analysis['image_size'] == 'small' else 4
        for j, element in enumerate(content_elements[:max_elements]):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            
            # Clean up the element text
            element_text = str(element).strip()
            if not element_text.startswith('•'):
                element_text = f"• {element_text}"
            
            p.text = element_text
            p.font.color.rgb = colors['text_dark']
            p.font.size = Pt(16 if content_analysis['image_size'] == 'large' else 14)
            p.space_before = Pt(8)
        
        # Add chart if slide has data and layout supports it
        if 'chart_data' in slide_data and content_analysis['recommended_layout'] != 'timeline':
            chart_x = 1.0 if content_analysis['image_size'] == 'small' else 1.5
            chart_width = 7.0 if content_analysis['image_size'] == 'small' else 6.0
            self.add_adaptive_chart(slide, slide_data['chart_data'], chart_x, chart_width, colors)
        
        # Enhanced footer with content-aware styling
        footer_text = f"AI-Generated Content | {content_analysis['content_type'].title()} Slide"
        self.create_enhanced_footer(slide, footer_text)

    def _create_cta_slide(self, prs, input_data: Dict):
        """Create call-to-action slide with dynamic styling."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Analyze CTA content for styling
        cta_text = input_data.get('strategic_plan', {}).get('primary_cta', 'Contact us to learn more!')
        cta_analysis = self.analyze_slide_content({'title': 'Next Steps', 'key_content_elements': [cta_text]})
        cta_analysis['action_oriented'] = True  # CTA slides are always action-oriented
        colors = self.get_dynamic_color_scheme('achievement')  # Use achievement colors for CTA
        
        # Dynamic CTA background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 90  # Vertical gradient for action
        fill.gradient_stops[0].color.rgb = colors['primary']
        fill.gradient_stops[1].color.rgb = colors['secondary']
        
        # Search for success/action background image
        cta_images = self.search_unsplash_image("success achievement handshake business") or self.get_fallback_images('success')
        if cta_images:
            self.add_background_image(slide, cta_images[0], opacity=0.25)
        
        # Action-oriented decorative elements
        self.add_smart_decorative_elements(slide, cta_analysis, colors)
        
        # Dynamic title
        title_shape = slide.shapes.title
        title_shape.text = "Next Steps"
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_shape.text_frame.paragraphs[0].font.name = 'Helvetica Neue'
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Enhanced CTA content with dynamic sizing
        cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2.5))
        tf = cta_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cta_text
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Dynamic font size based on CTA length
        cta_length = len(cta_text)
        if cta_length < 50:
            cta_font_size = 32
        elif cta_length < 100:
            cta_font_size = 26
        else:
            cta_font_size = 22
            
        p.font.size = Pt(cta_font_size)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add urgency indicators if CTA suggests immediate action
        if any(word in cta_text.lower() for word in ['now', 'today', 'immediate', 'start', 'book', 'schedule']):
            urgency_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2), Inches(0.8))
            urgency_tf = urgency_box.text_frame
            urgency_p = urgency_tf.paragraphs[0]
            urgency_p.text = "🚀 ACT NOW"
            urgency_p.font.color.rgb = colors['warning']
            urgency_p.font.size = Pt(14)
            urgency_p.font.bold = True
            urgency_p.alignment = PP_ALIGN.CENTER
        
        # Action-oriented footer
        footer_text = "Ready to Get Started | Take Action Today"
        self.create_enhanced_footer(slide, footer_text)

def main():
    """Main function to demonstrate the advanced PPT generator."""
    
    # Sample prompt and data
    sample_prompt = "Create a compelling sales presentation that highlights our AI solution's benefits with professional visuals and data-driven insights"
    
    sample_data = {
        "id": 67890,
        "slides": [
            {
                "id": 1,
                "title": "The Hidden Cost of Manual Sales Processes",
                "key_content_elements": [
                    "Sales reps spend 4-6 hours/week on manual tasks",
                    "Inconsistent documentation leads to lost insights",
                    "Manual processes reduce productivity by 30%",
                    "Revenue opportunities are missed due to inefficiency"
                ],
                "objective": "Highlight pain points and challenges",
                "chart_data": {
                    "categories": ["Manual Tasks", "Lost Insights", "Missed Opportunities"],
                    "values": [35, 25, 40]
                }
            },
            {
                "id": 2,
                "title": "Introducing GetAligned: Your AI Co-Pilot",
                "key_content_elements": [
                    "AI-powered insight extraction and automation",
                    "Seamless integration with existing tools",
                    "Real-time summary generation",
                    "Automated pitch deck creation"
                ],
                "objective": "Position our solution as the answer"
            }
        ],
        "strategic_plan": {
            "primary_cta": "Schedule a technical deep-dive with our team"
        },
        "title": "Unlock Sales Excellence: AI-Powered Insights for Success"
    }
    
    # Initialize generator
    generator = AdvancedPPTGenerator()
    
    # Generate presentation
    output_file = generator.generate_presentation_from_prompt(
        prompt=sample_prompt,
        input_data=sample_data,
        output_filename="AI_Enhanced_Sales_Presentation.pptx"
    )
    
    print(f"🎉 Advanced presentation generated: {output_file}")

if __name__ == "__main__":
    main()
